import io
import re
import zipfile
from collections import defaultdict
from pathlib import Path

import openpyxl
import docx
import pptx

from fastapi import HTTPException 

from typing import List, Union, Dict


def _parse_role_email_summary(role_email_summary: str) -> dict:
    if not role_email_summary or role_email_summary == "N/A":
        return {}

    email_data_by_role = {}
    
    # 按分号分割不同的角色部分
    role_sections = role_email_summary.split(';')

    for section in role_sections:
        # 分离角色名和成员列表
        if ':' not in section:
            continue
        role_key, _, member_list_str = section.partition(':')
        role_key = role_key.strip()
        
        if not role_key:
            continue

        output_lines = []
        # 按逗号分割每个成员
        members = member_list_str.strip().split(',')
        
        for member_str in members:
            member_str = member_str.strip()
            if not member_str:
                continue
            
            # 使用 rpartition 按最后一个下划线分割，以处理名字中可能包含下划线的情况
            display_name, separator, email = member_str.rpartition('_')
            
            if separator:  # 确保分割成功
                output_lines.append(display_name.strip())
                output_lines.append(email.strip())

        if output_lines:
            email_data_by_role[role_key] = "\n".join(output_lines)
            
    return email_data_by_role


def _prepare_profile_for_filling(profile_dict: dict) -> dict:
    # 1. 基本格式化：处理空值和列表
    formatted_profile = {}
    for key, value in profile_dict.items():
        if value is None or value == "" or value == []:
            formatted_profile[key] = "N/A"
        elif isinstance(value, list):
            formatted_profile[key] = ", ".join(map(str, value))
        else:
            formatted_profile[key] = str(value)

    # 2. 解析 role_summary 字符串 (用于姓名)
    role_summary_str = formatted_profile.get("role_summary", "N/A")
    role_data = {}
    if role_summary_str and role_summary_str != "N/A":
        pairs = role_summary_str.split(';')
        for pair in pairs:
            if ':' in pair:
                role, _, names = pair.partition(':')
                role_data[role.strip()] = names.strip()

    # 3. 定义占位符键与 role_summary 中的角色名称的映射
    # 键: 模板中的占位符 (如 <PMS.PJM>)
    # 值: role_summary 中的角色名称 (如 PjM)
    role_mapping = {
        'PJM': 'PjM',
        'TPM': 'TPM',
        'ECU_PCM': 'ECU-PCM',
        'SW_PCM': 'SW_PCM',
        'FSM': 'FSM',
        'SYS_ENG': 'Sys-ENG',
        'APP_PCM': 'App PCM',
        'HW_Dev': 'HW Developer',
        'AM': 'AM',
        'CM': 'CM',
        'COS': 'COS',
        'MECH_PCM': 'MECH-PCM',
        'SAMCO': 'SAMCO',
        'SEC': 'SEC',
        'TestM': 'Test Manager'
    }
    
    # 填充角色姓名字段
    for placeholder_key, summary_key in role_mapping.items():
        formatted_profile[placeholder_key] = role_data.get(summary_key, 'N/A')

    # 4. 创建一个反向映射，用于从源角色名找到占位符角色名
    # 例如：{'PjM': 'PJM', 'TPM': 'TPM', ...}
    inverted_role_mapping = {v: k for k, v in role_mapping.items()}

    # 5. 解析 role_email_summary 字符串
    email_summary_str = formatted_profile.get("role_email_summary", "N/A")
    parsed_emails_by_role = _parse_role_email_summary(email_summary_str)

    # 6. 根据反向映射，精确地转换邮件角色键
    corrected_email_data = {}
    for raw_key, value in parsed_emails_by_role.items():
        # 确保我们只处理以 'Email' 结尾的键
        if raw_key.endswith('Email'):
            # 提取前缀，例如从 'PjMEmail' 中提取 'PjM'
            prefix = raw_key[:-5]
            
            # 在反向映射中查找对应的大写占位符前缀
            placeholder_prefix = inverted_role_mapping.get(prefix)
            
            if placeholder_prefix:
                # 如果找到了，就构建新的、正确的键，例如 'PJM' + 'Email' -> 'PJMEmail'
                new_key = placeholder_prefix + 'Email'
                corrected_email_data[new_key] = value
            else:
                # 如果在映射中没找到，则保留原始键
                corrected_email_data[raw_key] = value
        else:
            # 如果键不是以 'Email' 结尾，也保留原始键
            corrected_email_data[raw_key] = value
            
    formatted_profile.update(corrected_email_data)

    return formatted_profile


def _is_macro_enabled(source: Union[Path, io.BytesIO]) -> bool:
    if isinstance(source, Path):
        return source.suffix.lower() == '.xlsm'

    if isinstance(source, io.BytesIO):
        try:
            source.seek(0)
            with zipfile.ZipFile(source, 'r') as zip_file:
                return 'xl/vbaProject.bin' in zip_file.namelist()
        except zipfile.BadZipFile:
            return False
        finally:
            source.seek(0)
    
    return False


def merge_and_fill_excel(profile_dict: dict, template_source: io.BytesIO, base_source: io.BytesIO) -> io.BytesIO:
    should_keep_vba = _is_macro_enabled(template_source)
    try:
        template_source.seek(0)
        base_source.seek(0)
        
        template_workbook = openpyxl.load_workbook(template_source, keep_vba=should_keep_vba)
        base_workbook = openpyxl.load_workbook(base_source, read_only=True)

        formatted_profile = _prepare_profile_for_filling(profile_dict)
        def replacer(match: re.Match) -> str:
            combined_key = match.group(1)
            if '-' in combined_key:
                keys = combined_key.split('-')
                value_parts = [formatted_profile.get(key, "N/A") for key in keys]
                return "_".join(map(str, value_parts))
            else:
                return formatted_profile.get(combined_key, "N/A")

        for template_sheet in template_workbook.worksheets:
            sheet_name = template_sheet.title
            
            if sheet_name in base_workbook.sheetnames:
                base_sheet = base_workbook[sheet_name]
                
                base_data = {
                    cell.coordinate: cell.value 
                    for row in base_sheet.iter_rows() 
                    for cell in row if cell.value is not None
                }

                template_coords = {
                    cell.coordinate 
                    for row in template_sheet.iter_rows() 
                    for cell in row if cell.value is not None
                }
                
                all_coords = set(base_data.keys()).union(template_coords)

                for coord in all_coords:
                    template_cell = template_sheet[coord]
                    original_template_value = template_cell.value

                    is_placeholder = isinstance(original_template_value, str) and "<PMS." in original_template_value
                    
                    if is_placeholder:
                        template_cell.value = re.sub(r"<PMS\.([^>]+)>", replacer, original_template_value)
                    
                    elif coord in base_data:
                        template_cell.value = base_data[coord]
                    

            else:
                print(f"警告: 模板中的工作表 '{sheet_name}' 在基础数据文件中未找到，将只处理占位符。")
                for row in template_sheet.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and "<PMS." in cell.value:
                            cell.value = re.sub(r"<PMS\.([^>]+)>", replacer, cell.value)

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=400, detail=f"处理Excel文件时出错: {e}")

    output_stream = io.BytesIO()
    template_workbook.save(output_stream)
    output_stream.seek(0)
    
    return output_stream


def fill_excel_by_placeholders(profile_dict: dict, source: Union[Path, io.BytesIO]) -> io.BytesIO:
    should_keep_vba = _is_macro_enabled(source)

    try:
        if isinstance(source, io.BytesIO):
            source.seek(0)
        workbook = openpyxl.load_workbook(source, keep_vba=should_keep_vba)
    except FileNotFoundError:
        raise HTTPException(status_code=500, detail=f"服务器错误: 模板文件未找到 '{source}'")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"无法处理提供的Excel文件: {e}")

    formatted_profile = _prepare_profile_for_filling(profile_dict)

    def replacer(match: re.Match) -> str:
        combined_key = match.group(1)
        if '-' in combined_key:
            keys = combined_key.split('-')
            value_parts = [formatted_profile.get(key, "N/A") for key in keys]
            return "_".join(value_parts)
        else:
            return formatted_profile.get(combined_key, "N/A")

    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and "<PMS." in cell.value:
                    cell.value = re.sub(r"<PMS\.([^>]+)>", replacer, cell.value)

    stream = io.BytesIO()
    workbook.save(stream)
    stream.seek(0)
    return stream


def fill_docx_by_placeholders(profile_dict: dict, source: Union[Path, io.BytesIO]) -> io.BytesIO:
    try:
        document = docx.Document(source)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"无法处理提供的Word文件: {e}")

    formatted_profile = _prepare_profile_for_filling(profile_dict)
    
    def replacer(match: re.Match) -> str:
        combined_key = match.group(1)
        if '-' in combined_key:
            keys = combined_key.split('-')
            value_parts = [formatted_profile.get(key, "N/A") for key in keys]
            return "_".join(value_parts)
        else:
            return formatted_profile.get(combined_key, "N/A")

    def substitute_in_paragraph(paragraph):
        if "<PMS." in paragraph.text:
            full_text = paragraph.text
            new_text = re.sub(r"<PMS\.([^>]+)>", replacer, full_text)
            if new_text != full_text:
                paragraph.text = new_text

    for paragraph in document.paragraphs:
        substitute_in_paragraph(paragraph)

    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    substitute_in_paragraph(paragraph)

    stream = io.BytesIO()
    document.save(stream)
    stream.seek(0)
    return stream


def fill_pptx_by_placeholders(profile_dict: dict, source: Union[Path, io.BytesIO]) -> io.BytesIO:
    try:
        presentation = pptx.Presentation(source)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"无法处理提供的PowerPoint文件: {e}")

    formatted_profile = _prepare_profile_for_filling(profile_dict)
    
    def replacer(match: re.Match) -> str:
        combined_key = match.group(1)
        if '-' in combined_key:
            keys = combined_key.split('-')
            value_parts = [formatted_profile.get(key, "N/A") for key in keys]
            return "_".join(value_parts)
        else:
            return formatted_profile.get(combined_key, "N/A")

    # 在PPTX中，换行符 \n 通常有效
    def substitute_in_text_frame(text_frame):
        for para in text_frame.paragraphs:
            if "<PMS." in para.text:
                full_text = para.text
                new_text = re.sub(r"<PMS\.([^>]+)>", replacer, full_text)
                if new_text != full_text:
                    para.text = new_text

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                substitute_in_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        substitute_in_text_frame(cell.text_frame)

    stream = io.BytesIO()
    presentation.save(stream)
    stream.seek(0)
    return stream
